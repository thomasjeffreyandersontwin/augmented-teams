#!/usr/bin/env python3
"""
Multi-format document parser for extracting text from various file types.
Supports: Word (.docx), Excel (.xlsx), PowerPoint (.pptx), PDF (.pdf), Markdown (.md), Text (.txt)

Libraries chosen are official/most-maintained parsers for each format.
Provides comprehensive coverage of business document types while maintaining a unified interface.
"""

from pathlib import Path
from typing import Dict, Optional
import logging

# Document parsers
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

logger = logging.getLogger(__name__)


class DocumentParser:
    """Extract text content from multiple document formats"""
    
    SUPPORTED_EXTENSIONS = {
        '.docx': 'word',
        '.xlsx': 'excel', 
        '.pptx': 'powerpoint',
        '.pdf': 'pdf',
        '.md': 'markdown',
        '.txt': 'text'
    }
    
    @classmethod
    def extract_text(cls, file_path: Path) -> Dict[str, any]:
        """
        Extract text from a document file.
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dict with keys: text, metadata, file_type
        """
        suffix = file_path.suffix.lower()
        
        if suffix not in cls.SUPPORTED_EXTENSIONS:
            logger.warning(f"Unsupported file type: {suffix} for {file_path}")
            return {
                'text': '',
                'metadata': {'error': f'Unsupported file type: {suffix}'},
                'file_type': 'unknown'
            }
        
        file_type = cls.SUPPORTED_EXTENSIONS[suffix]
        
        try:
            if file_type == 'word':
                return cls._extract_docx(file_path)
            elif file_type == 'excel':
                return cls._extract_xlsx(file_path)
            elif file_type == 'powerpoint':
                return cls._extract_pptx(file_path)
            elif file_type == 'pdf':
                return cls._extract_pdf(file_path)
            elif file_type == 'markdown':
                return cls._extract_markdown(file_path)
            elif file_type == 'text':
                return cls._extract_text(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return {
                'text': '',
                'metadata': {'error': str(e)},
                'file_type': file_type
            }
    
    @staticmethod
    def _extract_docx(file_path: Path) -> Dict:
        """Extract text from Word documents"""
        if Document is None:
            raise ImportError("python-docx not installed")
        
        doc = Document(file_path)
        
        # Extract paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        # Extract tables
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text:
                    tables_text.append(row_text)
        
        full_text = '\n\n'.join(paragraphs)
        if tables_text:
            full_text += '\n\nTables:\n' + '\n'.join(tables_text)
        
        return {
            'text': full_text,
            'metadata': {
                'paragraphs': len(paragraphs),
                'tables': len(doc.tables),
                'chars': len(full_text)
            },
            'file_type': 'word'
        }
    
    @staticmethod
    def _extract_xlsx(file_path: Path) -> Dict:
        """Extract text from Excel spreadsheets"""
        if load_workbook is None:
            raise ImportError("openpyxl not installed")
        
        wb = load_workbook(file_path, data_only=True, read_only=True)
        
        all_text = []
        sheet_metadata = {}
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            # Extract all cells with values
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                if row_text:
                    sheet_text.append(row_text)
            
            all_text.extend(sheet_text)
            sheet_metadata[sheet_name] = len(sheet_text) - 1  # Exclude header
        
        wb.close()
        full_text = '\n'.join(all_text)
        
        return {
            'text': full_text,
            'metadata': {
                'sheets': list(wb.sheetnames),
                'sheet_rows': sheet_metadata,
                'chars': len(full_text)
            },
            'file_type': 'excel'
        }
    
    @staticmethod
    def _extract_pptx(file_path: Path) -> Dict:
        """Extract text from PowerPoint presentations"""
        if Presentation is None:
            raise ImportError("python-pptx not installed")
        
        prs = Presentation(file_path)
        
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {i}:"]
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # More than just the header
                slides_text.append('\n'.join(slide_text))
        
        full_text = '\n\n'.join(slides_text)
        
        return {
            'text': full_text,
            'metadata': {
                'slides': len(prs.slides),
                'chars': len(full_text)
            },
            'file_type': 'powerpoint'
        }
    
    @staticmethod
    def _extract_pdf(file_path: Path) -> Dict:
        """Extract text from PDF files"""
        if PdfReader is None:
            raise ImportError("pypdf not installed")
        
        reader = PdfReader(file_path)
        
        pages_text = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text.strip():
                pages_text.append(f"Page {i}:\n{text}")
        
        full_text = '\n\n'.join(pages_text)
        
        return {
            'text': full_text,
            'metadata': {
                'pages': len(reader.pages),
                'chars': len(full_text)
            },
            'file_type': 'pdf'
        }
    
    @staticmethod
    def _extract_markdown(file_path: Path) -> Dict:
        """Extract text from Markdown files"""
        content = file_path.read_text(encoding='utf-8')
        
        # Keep markdown formatting for better semantic understanding
        return {
            'text': content,
            'metadata': {
                'lines': len(content.split('\n')),
                'chars': len(content)
            },
            'file_type': 'markdown'
        }
    
    @staticmethod
    def _extract_text(file_path: Path) -> Dict:
        """Extract text from plain text files"""
        content = file_path.read_text(encoding='utf-8')
        
        return {
            'text': content,
            'metadata': {
                'lines': len(content.split('\n')),
                'chars': len(content)
            },
            'file_type': 'text'
        }


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python document_parsers.py <file_path>")
        sys.exit(1)
    
    file_path = Path(sys.argv[1])
    
    if not file_path.exists():
        print(f"File not found: {file_path}")
        sys.exit(1)
    
    result = DocumentParser.extract_text(file_path)
    
    print(f"File Type: {result['file_type']}")
    print(f"Metadata: {result['metadata']}")
    print(f"\nExtracted Text ({len(result['text'])} chars):")
    print("-" * 60)
    print(result['text'][:500])  # First 500 chars
    if len(result['text']) > 500:
        print("...")

