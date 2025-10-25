#!/usr/bin/env python3
"""
Git Integration FastAPI Service

A FastAPI wrapper around git_integration.py functions.
Provides REST endpoints for all Git operations.
"""

import os
import sys
from pathlib import Path
from typing import Optional
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Header, Depends
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn
import requests
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# GitHub authentication
GITHUB_TOKEN = os.getenv("GITHUB_TOKEN")
GITHUB_REPO = os.getenv("GITHUB_REPO", "thomasjeffreyandersontwin/augmented-teams")
GITHUB_API_BASE = "https://api.github.com"

# Use GitHub token for client authentication too
CLIENT_TOKEN = GITHUB_TOKEN

# Security scheme
security = HTTPBearer(auto_error=False)

# Add the current directory to Python path to import git_integration
current_dir = Path(__file__).parent
sys.path.insert(0, str(current_dir))

# Import our git integration functions
from integration import (
    ensure_latest, save_code, commit_and_push, copy_workflow_files,
    get_repository_tree, get_folder_content, search_git_content, 
    search_git_files, get_file_content, run_cmd, REPO_PATH
)

def get_github_headers():
    """Get GitHub API headers with authentication"""
    headers = {
        "Accept": "application/vnd.github.v3+json",
        "User-Agent": "Git-Integration-Service"
    }
    if GITHUB_TOKEN:
        headers["Authorization"] = f"token {GITHUB_TOKEN}"
    return headers

def check_github_auth():
    """Check if GitHub authentication is working"""
    if not GITHUB_TOKEN:
        return False, "GITHUB_TOKEN not set in environment"
    
    try:
        response = requests.get(f"{GITHUB_API_BASE}/user", headers=get_github_headers())
        if response.status_code == 200:
            user_data = response.json()
            return True, f"Authenticated as {user_data.get('login', 'unknown')}"
        else:
            return False, f"GitHub auth failed: {response.status_code}"
    except Exception as e:
        return False, f"GitHub auth error: {e}"

def verify_client_token(credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Verify client authentication token - MANDATORY"""
    if not CLIENT_TOKEN:
        raise HTTPException(
            status_code=500, 
            detail="Service not configured: GITHUB_TOKEN required"
        )
    
    if not credentials:
        raise HTTPException(
            status_code=401, 
            detail="Authentication required. Provide token in Authorization header."
        )
    
    if credentials.credentials != CLIENT_TOKEN:
        raise HTTPException(
            status_code=403, 
            detail="Invalid authentication token"
        )
    
    return True

def get_current_client():
    """Get current authenticated client info"""
    return {
        "authenticated": True,
        "service": "git-integration",
        "version": "1.0.0"
    }

app = FastAPI(
    title="Git Integration Service",
    description="REST API for Git operations using native git commands",
    version="1.0.0"
)

# Add CORS middleware for GPT Actions
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for GPT Actions
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models
class CommitRequest(BaseModel):
    content: str
    file_path: str
    commit_message: Optional[str] = None

class CommitResponse(BaseModel):
    success: bool
    message: str
    file_path: str

class SearchRequest(BaseModel):
    query: str
    file_pattern: Optional[str] = "*"

@app.get("/")
async def health_check():
    """Health check endpoint with GitHub auth status"""
    auth_status, auth_message = check_github_auth()
    return {
        "message": "Git Integration Service is running",
        "service": "git-integration",
        "version": "1.0.0",
        "github_auth": auth_status,
        "github_message": auth_message,
        "repository": GITHUB_REPO
    }

@app.get("/auth/github")
async def github_auth_status():
    """Check GitHub authentication status"""
    auth_status, auth_message = check_github_auth()
    return {
        "authenticated": auth_status,
        "message": auth_message,
        "repository": GITHUB_REPO,
        "token_set": bool(GITHUB_TOKEN)
    }

@app.get("/auth/client")
async def client_auth_status():
    """Check client authentication configuration"""
    return {
        "authentication_required": True,
        "token_configured": bool(CLIENT_TOKEN),
        "auth_method": "Bearer token in Authorization header (GITHUB_TOKEN)"
    }

@app.get("/auth/verify")
async def verify_auth(client: dict = Depends(verify_client_token)):
    """Verify client authentication"""
    return {
        "authenticated": True,
        "message": "Authentication successful",
        "client": get_current_client()
    }

@app.get("/status")
async def get_status():
    """Get repository status"""
    try:
        stdout, stderr, returncode = run_cmd(["git", "status", "--porcelain"])
        return {
            "has_changes": bool(stdout),
            "changes": stdout.split('\n') if stdout else [],
            "error": stderr if returncode != 0 else None
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/commit/text", response_model=CommitResponse)
async def commit_text(request: CommitRequest):
    """Commit text content to a file"""
    try:
        # Save the content
        file_path = save_code(request.content, request.file_path.split('/')[-1], 
                            '/'.join(request.file_path.split('/')[:-1]))
        
        # Commit and push
        commit_and_push(request.commit_message)
        
        return CommitResponse(
            success=True,
            message=f"Successfully committed {request.file_path}",
            file_path=request.file_path
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/commit/document", response_model=CommitResponse)
async def commit_document(
    file: UploadFile = File(...),
    file_path: str = Form(...),
    commit_message: Optional[str] = Form(None)
):
    """Commit uploaded document"""
    try:
        # Read uploaded file content
        content = await file.read()
        
        # Save as text file (for now)
        text_content = content.decode('utf-8') if isinstance(content, bytes) else str(content)
        save_code(text_content, file_path.split('/')[-1], 
                 '/'.join(file_path.split('/')[:-1]))
        
        # Commit and push
        commit_and_push(commit_message)
        
        return CommitResponse(
            success=True,
            message=f"Successfully committed {file_path}",
            file_path=file_path
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/push")
async def push_changes():
    """Push committed changes to remote"""
    try:
        commit_and_push()
        return {
            "success": True,
            "message": "Successfully pushed changes to remote"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/tree")
async def get_tree(path: str = "."):
    """Get repository tree structure"""
    try:
        result = get_repository_tree(path)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/folder/{folder_path:path}")
async def get_folder(folder_path: str):
    """Get content of a specific folder"""
    try:
        result = get_folder_content(folder_path)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/search/content")
async def search_content(request: SearchRequest):
    """Search repository content using git grep"""
    try:
        result = search_git_content(request.query, request.file_pattern)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/search/files")
async def search_files(pattern: str):
    """Search for files by name pattern"""
    try:
        result = search_git_files(pattern)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/file/{file_path:path}")
async def get_file(file_path: str):
    """Get content of a specific file"""
    try:
        result = get_file_content(file_path)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/sync")
async def sync_repository():
    """Sync repository (pull latest changes)"""
    try:
        ensure_latest()
        return {
            "success": True,
            "message": "Repository synced successfully"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/delete/{file_path:path}", dependencies=[Depends(verify_client_token)])
async def delete_file(file_path: str, commit_message: Optional[str] = None):
    """Delete a file from the repository"""
    try:
        # Check if file exists
        full_path = REPO_PATH / file_path
        print(f"DEBUG: Looking for file at {full_path}")
        if not full_path.exists():
            raise HTTPException(status_code=404, detail=f"File {file_path} not found at {full_path}")
        
        # Remove the file
        stdout, stderr, returncode = run_cmd(["git", "rm", file_path])
        if returncode != 0:
            raise HTTPException(status_code=500, detail=f"git rm failed: {stderr}")
        
        # Commit the deletion
        msg = commit_message or f"Remove {file_path}"
        stdout, stderr, returncode = run_cmd(["git", "commit", "-m", msg])
        if returncode != 0:
            raise HTTPException(status_code=500, detail=f"git commit failed: {stderr}")
        
        # Push changes
        stdout, stderr, returncode = run_cmd(["git", "push", "origin", "main"])
        if returncode != 0:
            raise HTTPException(status_code=500, detail=f"git push failed: {stderr}")
        
        return {
            "success": True,
            "message": f"Successfully deleted {file_path}",
            "file_path": file_path
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/workflows/copy")
async def copy_workflows():
    """Copy workflow files to .github/workflows/"""
    try:
        copy_workflow_files()
        return {
            "success": True,
            "message": "Workflow files copied successfully"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class DeleteFileRequest(BaseModel):
    file_path: str
    commit_message: Optional[str] = None

@app.post("/deleteFile", dependencies=[Depends(verify_client_token)])
async def delete_file_post(request: DeleteFileRequest):
    """Delete a file from the repository (POST endpoint for GPT Actions compatibility)"""
    try:
        # Check if file exists
        full_path = REPO_PATH / request.file_path
        if not full_path.exists():
            raise HTTPException(status_code=404, detail=f"File {request.file_path} not found at {full_path}")
        
        # Remove the file
        stdout, stderr, returncode = run_cmd(["git", "rm", request.file_path])
        if returncode != 0:
            raise HTTPException(status_code=500, detail=f"git rm failed: {stderr}")
        
        # Commit the deletion
        msg = request.commit_message or f"Remove {request.file_path}"
        stdout, stderr, returncode = run_cmd(["git", "commit", "-m", msg])
        if returncode != 0:
            raise HTTPException(status_code=500, detail=f"git commit failed: {stderr}")
        
        # Push changes
        stdout, stderr, returncode = run_cmd(["git", "push", "origin", "main"])
        if returncode != 0:
            raise HTTPException(status_code=500, detail=f"git push failed: {stderr}")
        
        return {
            "success": True,
            "message": f"Successfully deleted {request.file_path}",
            "file_path": request.file_path
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def extract_text_from_binary(binary_data: bytes, file_extension: str) -> str:
    """Extract text from binary files"""
    try:
        if file_extension.lower() == ".pptx":
            from pptx import Presentation
            from io import BytesIO
            
            prs = Presentation(BytesIO(binary_data))
            slides_text = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
            return "\n\n".join(slides_text)
        
        elif file_extension.lower() == ".pdf":
            try:
                import PyPDF2
                from io import BytesIO
                
                pdf_reader = PyPDF2.PdfReader(BytesIO(binary_data))
                text_content = []
                for page_num, page in enumerate(pdf_reader.pages, start=1):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"--- Page {page_num} ---\n{page_text}")
                return "\n\n".join(text_content)
            except ImportError:
                return "[PDF extraction requires PyPDF2: pip install PyPDF2]"
        
        elif file_extension.lower() in [".docx", ".doc"]:
            try:
                from docx import Document
                from io import BytesIO
                
                doc = Document(BytesIO(binary_data))
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text.strip())
                return "\n".join(paragraphs)
            except ImportError:
                return "[DOCX extraction requires python-docx: pip install python-docx]"
        
        else:
            return f"[Unsupported binary file type: {file_extension}]"
    
    except Exception as e:
        return f"[Error extracting text from {file_extension}: {str(e)}]"

@app.get("/extract/{file_path:path}", dependencies=[Depends(verify_client_token)])
async def extract_file_content(file_path: str):
    """Extract text content from any file in the repository (including binary files)"""
    try:
        import base64
        
        # Check if file exists
        full_path = REPO_PATH / file_path
        if not full_path.exists():
            raise HTTPException(status_code=404, detail=f"File {file_path} not found at {full_path}. REPO_PATH: {REPO_PATH}")
        
        # Get file extension and name
        file_extension = full_path.suffix.lower()
        file_name = full_path.name.lower()
        
        # Read file content
        with open(full_path, "rb") as f:
            file_content = f.read()
        
        # Known text file extensions
        text_extensions = [".txt", ".md", ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".xml", ".html", ".css", ".sh", ".bash", ".ps1", ".sql", ".log", ".ini", ".cfg", ".conf"]
        
        # Known text files without extensions
        text_filenames = [".gitignore", ".gitattributes", "readme", "license", "dockerfile", "makefile", ".env", ".dockerignore"]
        
        # Check if it's a known text file
        is_text_file = (file_extension in text_extensions) or (file_name in text_filenames) or (file_extension == "" and not file_name.startswith("."))
        
        # Try to decode as text first (for text files or files without extension)
        if is_text_file or file_extension == "":
            try:
                text_content = file_content.decode('utf-8')
                return {
                    "success": True,
                    "file_path": file_path,
                    "file_type": "text",
                    "content": text_content,
                    "size": len(file_content)
                }
            except UnicodeDecodeError:
                # If decode fails for a known text file, it's an error
                if is_text_file:
                    return {
                        "success": False,
                        "file_path": file_path,
                        "error": "File contains non-UTF-8 text content"
                    }
                # Otherwise, treat as binary and continue
        
        # It's a binary file, try to extract text
        extracted_text = extract_text_from_binary(file_content, file_extension)
        return {
            "success": True,
            "file_path": file_path,
            "file_type": "binary",
            "extracted_content": extracted_text,
            "size": len(file_content),
            "extension": file_extension
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    print("Starting Git Integration Service on port 8001...")
    print(f"REPO_PATH: {REPO_PATH}")
    print(f"REPO_PATH exists: {REPO_PATH.exists()}")
    print(f"Current working directory: {Path.cwd()}")
    uvicorn.run(app, host="0.0.0.0", port=8001, log_level="info")

