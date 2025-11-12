from pptx import Presentation

prs = Presentation('1) Story Writing Training.pptx')

slides_data = []

for slide_num, slide in enumerate(prs.slides, 1):
    slide_info = {
        'slide_number': slide_num,
        'title': '',
        'text_content': [],
        'shapes_info': [],
        'tables': [],
        'images': []
    }
    
    # Extract title
    if slide.shapes.title:
        slide_info['title'] = slide.shapes.title.text
    
    # Extract all text and shape info
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text and text != slide_info['title']:
                    slide_info['text_content'].append(text)
        
        # Check if shape contains a table
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            slide_info['tables'].append(table_data)
        
        # Track images
        if hasattr(shape, 'image'):
            slide_info['images'].append({'type': 'image', 'name': shape.name})
        
        # Track other shapes (diagrams, charts, etc)
        slide_info['shapes_info'].append({
            'name': shape.name,
            'type': str(shape.shape_type),
            'has_text': shape.has_text_frame,
            'has_table': shape.has_table
        })
    
    slides_data.append(slide_info)

# Print structured output
for slide in slides_data:
    print(f'\n{"="*80}')
    print(f'SLIDE {slide["slide_number"]}: {slide["title"]}')
    print(f'{"="*80}')
    
    if slide['text_content']:
        print('\n[Text Content]:')
        for text in slide['text_content']:
            print(f'  - {text}')
    
    if slide['tables']:
        print('\n[Tables]:')
        for table_idx, table in enumerate(slide['tables'], 1):
            print(f'\n  Table {table_idx}:')
            for row_idx, row in enumerate(table):
                print(f'  | {" | ".join(row)} |')
    
    if slide['images']:
        print(f'\n[Images]: {len(slide["images"])} found')
        for img in slide['images']:
            print(f'  - {img["name"]}')
    
    # Show diagram/shape info that might contain visual examples
    non_text_shapes = [s for s in slide['shapes_info'] if not s['has_text'] and not s['has_table']]
    if non_text_shapes:
        print(f'\n[Visual Shapes/Diagrams]: {len(non_text_shapes)} found')
        for shape in non_text_shapes:
            print(f'  - {shape["name"]} (Type: {shape["type"]})')

print('\n' + '='*80)
print('EXTRACTION COMPLETE')
print('='*80)

